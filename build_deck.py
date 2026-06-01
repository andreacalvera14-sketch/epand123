#!/usr/bin/env python3
"""
build_deck.py — Epsilon CSA Deck builder (python-pptx).

Generates a fully editable 16:9 PowerPoint (Epsilon_CSA_Deck.pptx)
styled in Epsilon's black / white / turquoise brand identity.

Colour palette
--------------
  BLACK  #0B0B0B  — headings, body text
  WHITE  #FFFFFF  — slide background
  TURQ   #00B0A3  — Publicis Turquoise accent (rules, callout bars only)
  GREY   #6B6B6B  — secondary text, captions, footers
  LIGHT  #E7E7E7  — table borders

Font: Helvetica Neue (falls back to Arial on Windows)

Content source: epsilon-deck.md (Marp reference kept in the repo).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ---------------------------------------------------------------------------
# Brand constants
# ---------------------------------------------------------------------------
BLACK = RGBColor(0x0B, 0x0B, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TURQ  = RGBColor(0x00, 0xB0, 0xA3)   # Publicis Turquoise — accent only
GREY  = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT = RGBColor(0xE7, 0xE7, 0xE7)   # for table borders
FONT  = "Helvetica Neue"

SLIDE_W      = Inches(13.333)
SLIDE_H      = Inches(7.5)
TOTAL_SLIDES = 8
FOOTER_TEXT  = (
    "[Your Name]  \u00b7  Client Success Associate"
    "  \u00b7  Epsilon London  \u00b7  June 2026"
)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def new_prs():
    """Create a blank 16:9 presentation."""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    """Add and return a completely blank slide (layout 6 = blank)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fmt_run(run, font_name, font_size, bold, italic, color):
    """Apply formatting to a single text run."""
    run.font.name   = font_name
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, text, x, y, w, h,
                font_name=FONT, font_size=14, bold=False, italic=False,
                color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    """Add a single-run text box and return the shape."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _fmt_run(run, font_name, font_size, bold, italic, color)
    return txBox


def add_para(tf, text, font_name=FONT, font_size=14, bold=False, italic=False,
             color=BLACK, align=PP_ALIGN.LEFT, space_before=Pt(4)):
    """Append a paragraph (single run) to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    _fmt_run(run, font_name, font_size, bold, italic, color)
    return p


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    """
    Add a filled rectangle shape with no visible border (or a specified one).
    Used for accent rules and callout bars.
    """
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()   # transparent line
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    return shape


# ---------------------------------------------------------------------------
# Slide-level helpers
# ---------------------------------------------------------------------------

def add_accent_rule(slide, x, y, width=Inches(1.5)):
    """Add the thin 1pt TURQUOISE horizontal rule that sits under each title."""
    return add_rect(slide, x, y, width, Pt(1.5), TURQ)


def add_slide_header(slide, title, slide_num, subtitle=None):
    """
    Render the standard non-title header:
      - 32pt bold BLACK title top-left
      - Thin TURQUOISE rule below title
      - Optional 16pt GREY subtitle under rule
      - 10pt GREY page-number top-right  (e.g. "02 / 08")
      - 9pt GREY footer bottom-left

    Returns the y-coordinate just below the header, ready for slide content.
    """
    # Title
    add_textbox(slide, title,
                x=Inches(0.5), y=Inches(0.28),
                w=Inches(10), h=Inches(0.6),
                font_size=32, bold=True, color=BLACK)

    # Accent rule under title
    add_accent_rule(slide, x=Inches(0.5), y=Inches(0.95))

    # Optional subtitle
    if subtitle:
        add_textbox(slide, subtitle,
                    x=Inches(0.5), y=Inches(1.05),
                    w=Inches(10), h=Inches(0.38),
                    font_size=16, color=GREY)

    # Page number — top-right
    pg_text = f"{slide_num:02d} / {TOTAL_SLIDES:02d}"
    add_textbox(slide, pg_text,
                x=Inches(12.1), y=Inches(0.22),
                w=Inches(1.0), h=Inches(0.3),
                font_size=10, color=GREY, align=PP_ALIGN.RIGHT)

    # Footer — bottom-left
    add_textbox(slide, FOOTER_TEXT,
                x=Inches(0.5), y=Inches(7.12),
                w=Inches(11), h=Inches(0.28),
                font_size=9, color=GREY)

    return Inches(1.55) if subtitle else Inches(1.2)


def add_callout(slide, text, x, y, w, h):
    """
    Render a quote/callout block:
      - 2pt wide TURQUOISE vertical bar on the left
      - 12pt italic GREY text to the right of the bar
    """
    bar_w = Pt(3)
    # Vertical TURQUOISE bar
    add_rect(slide, x, y, bar_w, h, TURQ)
    # Quote text, indented from bar
    txBox = slide.shapes.add_textbox(x + bar_w + Pt(4), y, w - bar_w - Pt(4), h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    _fmt_run(run, FONT, 12, False, True, GREY)
    return txBox


# ---------------------------------------------------------------------------
# Table helpers
# ---------------------------------------------------------------------------

def set_cell_text(cell, text, font_size=12, bold=False, color=BLACK):
    """Clear a table cell and write formatted text."""
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    _fmt_run(run, FONT, font_size, bold, False, color)


def style_header_row(table):
    """BLACK fill + WHITE bold text for the header row."""
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = WHITE
                run.font.bold = True


def style_body_rows(table):
    """WHITE fill for all non-header rows."""
    for row_idx, row in enumerate(table.rows):
        if row_idx == 0:
            continue
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    """
    Title slide — white background, large bold heading, turquoise accent,
    pull-quote, tag line.
    """
    slide = blank_slide(prs)

    # Main heading
    add_textbox(slide, "Client Success Associate",
                x=Inches(1.2), y=Inches(2.0),
                w=Inches(11.0), h=Inches(1.0),
                font_size=54, bold=True, color=BLACK,
                align=PP_ALIGN.CENTER)

    # Sub-heading
    add_textbox(slide, "Supporting a Live Retail Media Account",
                x=Inches(1.2), y=Inches(3.1),
                w=Inches(11.0), h=Inches(0.6),
                font_size=28, color=BLACK,
                align=PP_ALIGN.CENTER)

    # Short turquoise accent rule (1" wide, centred)
    rule_w = Inches(1.0)
    add_accent_rule(slide,
                    x=(SLIDE_W - rule_w) / 2,
                    y=Inches(3.85),
                    width=rule_w)

    # "[Your Name]  ·  Final round  ·  Epsilon London  ·  June 2026"
    add_textbox(slide,
                "[Your Name]  \u00b7  Final round  \u00b7  Epsilon London  \u00b7  June 2026",
                x=Inches(1.2), y=Inches(4.05),
                w=Inches(11.0), h=Inches(0.4),
                font_size=14, color=GREY,
                align=PP_ALIGN.CENTER)

    # Pull-quote — lower third
    add_textbox(slide,
                "\u201cKeep delivery on track, make reporting useful,"
                " and make the Client Lead look great.\u201d",
                x=Inches(2.0), y=Inches(5.0),
                w=Inches(9.3), h=Inches(0.75),
                font_size=18, italic=True, color=BLACK,
                align=PP_ALIGN.CENTER)

    # Tag line
    add_textbox(slide,
                "Organise   \u00b7   Analyse   \u00b7   Communicate",
                x=Inches(1.2), y=Inches(6.2),
                w=Inches(11.0), h=Inches(0.4),
                font_size=14, color=GREY,
                align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 2 — Operational Management
# ---------------------------------------------------------------------------

def slide_02_operations(prs):
    """Two-column layout: My operating system | Prioritisation rule."""
    slide = blank_slide(prs)
    content_y = add_slide_header(slide,
                                 "Operational Management",
                                 slide_num=2,
                                 subtitle="Staying organised across workstreams")

    # ── LEFT column ─────────────────────────────────────────────────────────
    lx = Inches(0.5)
    col_w = Inches(5.9)

    add_textbox(slide, "My operating system",
                x=lx, y=content_y,
                w=col_w, h=Inches(0.38),
                font_size=14, bold=True, color=BLACK)

    os_bullets = [
        "— Single source of truth tracker per client: Workstream \u00b7 Owner"
        " \u00b7 Status \u00b7 Next milestone \u00b7 R/A/G \u00b7 Last client update",
        "— Daily 15-min self stand-up: must-dos vs nice-to-dos",
        "— Friday wrap note to Client Lead: shipped / slipped / coming",
        "— Shared launch calendar synced to retailer trading calendar",
    ]
    txBox = slide.shapes.add_textbox(lx, content_y + Inches(0.45), col_w, Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(os_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = item
        _fmt_run(run, FONT, 13, False, False, BLACK)

    # ── Vertical divider ────────────────────────────────────────────────────
    add_rect(slide,
             x=Inches(6.6), y=content_y,
             w=Pt(1), h=Inches(4.5),
             fill_color=LIGHT)

    # ── RIGHT column ────────────────────────────────────────────────────────
    rx = Inches(7.0)

    add_textbox(slide, "Prioritisation rule",
                x=rx, y=content_y,
                w=col_w, h=Inches(0.38),
                font_size=14, bold=True, color=BLACK)

    pr_bullets = [
        "— MoSCoW + deadline proximity: client-facing this week wins",
        "— \u201cTouch it once\u201d for inbox & Slack",
        "— Always know the next 3 deadlines without looking",
    ]
    txBox2 = slide.shapes.add_textbox(rx, content_y + Inches(0.45), col_w, Inches(2.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(pr_bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = item
        _fmt_run(run, FONT, 13, False, False, BLACK)


# ---------------------------------------------------------------------------
# Slide 3 — Reporting & QA
# ---------------------------------------------------------------------------

def slide_03_reporting(prs):
    """Cadence table + QA checklist + insight callout."""
    slide = blank_slide(prs)
    content_y = add_slide_header(slide,
                                 "Reporting & QA",
                                 slide_num=3,
                                 subtitle="Turning data into action")

    # ── Cadence table ───────────────────────────────────────────────────────
    tbl_x = Inches(0.5)
    tbl_w = Inches(8.3)
    table = slide.shapes.add_table(
        4, 3, tbl_x, content_y, tbl_w, Inches(1.5)
    ).table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(2.7)
    table.columns[2].width = Inches(4.1)

    for ci, hdr in enumerate(["Cadence", "Audience", "Focus"]):
        set_cell_text(table.cell(0, ci), hdr, font_size=13, bold=True, color=WHITE)

    rows_data = [
        ("Weekly",  "Brand / agency",
         "Pacing, ROAS, CTR/CVR, top/bottom SKUs, 1\u20132 actions"),
        ("Monthly", "Client Lead + brand marketing",
         "Trend, SoV, learnings, next month\u2019s plan"),
        ("QBR",     "Senior stakeholders",
         "Strategic story, incrementality, roadmap"),
    ]
    for ri, (c0, c1, c2) in enumerate(rows_data):
        set_cell_text(table.cell(ri + 1, 0), c0, font_size=12)
        set_cell_text(table.cell(ri + 1, 1), c1, font_size=12)
        set_cell_text(table.cell(ri + 1, 2), c2, font_size=12)

    style_header_row(table)
    style_body_rows(table)

    # ── QA checklist ────────────────────────────────────────────────────────
    qa_y = content_y + Inches(1.65)
    add_textbox(slide, "QA checklist \u2014 every client deliverable",
                x=Inches(0.5), y=qa_y,
                w=Inches(8.3), h=Inches(0.35),
                font_size=13, bold=True, color=BLACK)

    qa_items = [
        "1.  Numbers reconcile: platform \u2194 report \u2194 invoice",
        "2.  Date range, currency, time zone correct",
        "3.  Pacing % matches days elapsed",
        "4.  Spelling, client name, logos, tone",
        "5.  Every chart has a \u201cso what?\u201d",
        "6.  Peer review \u2014 fresh eyes catch tired-eye mistakes",
    ]
    txBox = slide.shapes.add_textbox(
        Inches(0.5), qa_y + Inches(0.42), Inches(8.3), Inches(2.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(qa_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = item
        _fmt_run(run, FONT, 12, False, False, BLACK)

    # ── Insight callout — right side ────────────────────────────────────────
    callout_text = (
        'Before: "CTR is 0.42%"\n\n'
        'After: "CTR on Brand A sponsored products \u221218% WoW, driven by 3 SKUs going OOS. '
        'Recommend pausing those line items and reallocating \u00a34k to top-performing Kids."'
    )
    add_callout(slide, callout_text,
                x=Inches(9.1), y=content_y,
                w=Inches(3.9), h=Inches(4.2))


# ---------------------------------------------------------------------------
# Slide 4 — Issue Triage & Escalation
# ---------------------------------------------------------------------------

def slide_04_triage(prs):
    """Scenario + 6-step triage + escalation script callout."""
    slide = blank_slide(prs)
    content_y = add_slide_header(slide,
                                 "Issue Triage & Escalation",
                                 slide_num=4,
                                 subtitle="Worked example")

    # Scenario line
    add_textbox(slide,
                "Scenario: Friday 4pm \u2014 \u201cSponsored products spend dropped 60% overnight\u201d",
                x=Inches(0.5), y=content_y,
                w=Inches(12.3), h=Inches(0.38),
                font_size=13, italic=True, color=GREY)

    # 6-step triage
    steps = [
        "1.  Acknowledge in 15 min \u2014 \u201cInvestigating, will revert by 5:30pm.\u201d",
        "2.  Diagnose root cause \u2014 platform \u00b7 pacing/caps \u00b7 bids \u00b7 SKU feed \u00b7 taxonomy",
        "3.  Confirm with data, not assumption",
        "4.  Escalate to Client Lead \u2014 options + recommendation",
        "5.  Reply to client \u2014 cause, plan, ETA",
        "6.  Monday post-mortem \u2014 add feed-health alert",
    ]
    txBox = slide.shapes.add_textbox(
        Inches(0.5), content_y + Inches(0.5), Inches(8.0), Inches(3.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = step
        _fmt_run(run, FONT, 13, False, False, BLACK)

    # Escalation script — callout right side
    esc_text = (
        "Root cause: 40% of eligible SKUs OOS via 2am feed sync failure.\n\n"
        "(A) Hold spend, fix feed, relaunch Mon \u2014 ~\u00a38k underspend.\n"
        "(B) Reallocate to top 20 in-stock SKUs today \u2014 protects pacing.\n\n"
        "Recommend B + feed-health alert. Need decision by 5pm."
    )
    add_callout(slide, esc_text,
                x=Inches(8.8), y=content_y + Inches(0.45),
                w=Inches(4.3), h=Inches(3.5))


# ---------------------------------------------------------------------------
# Slide 5 — Optimisation Examples
# ---------------------------------------------------------------------------

def slide_05_optimisation(prs):
    """3-column Did / Why / Impact table."""
    slide = blank_slide(prs)
    content_y = add_slide_header(slide,
                                 "Optimisation Examples",
                                 slide_num=5,
                                 subtitle="Did \u2192 Why \u2192 Impact")

    # Table: header + 3 example rows
    table = slide.shapes.add_table(
        4, 3, Inches(0.5), content_y, Inches(12.3), Inches(4.8)
    ).table
    for ci in range(3):
        table.columns[ci].width = Inches(4.1)

    for ci, hdr in enumerate(["Did", "Why", "Impact"]):
        set_cell_text(table.cell(0, ci), hdr, font_size=13, bold=True, color=WHITE)

    examples = [
        (
            "SKU-level bid re-allocation: rebalanced bids, paused long-tail"
            " <0.1% CTR / 14 days",
            "Pareto \u2014 long-tail ate budget without converting",
            "ROAS +32%, sales +18%, spend flat",
        ),
        (
            "Search-term harvesting: promoted converting non-brand terms to"
            " manual; losers \u2192 negatives",
            "Capture incremental non-brand demand at controlled CPC",
            "New-to-brand sales +24%, wasted spend \u2212\u00a33.2k/mo",
        ),
        (
            "Creative refresh: A/B tested 3 creatives for Back-to-School"
            " after CTR fatigue",
            "Align to retailer\u2019s trading moments to lift SoV",
            "CTR +41%, target hit 9 days early",
        ),
    ]
    for ri, (did, why, impact) in enumerate(examples):
        set_cell_text(table.cell(ri + 1, 0), did,  font_size=12)
        set_cell_text(table.cell(ri + 1, 1), why,  font_size=12)
        set_cell_text(table.cell(ri + 1, 2), impact, font_size=12, bold=True)

    style_header_row(table)
    style_body_rows(table)

    # Footnote
    add_textbox(slide, "(figures illustrative)",
                x=Inches(0.5), y=Inches(6.8),
                w=Inches(5), h=Inches(0.28),
                font_size=10, italic=True, color=GREY)


# ---------------------------------------------------------------------------
# Slide 6 — Client Comms & QBR Readiness
# ---------------------------------------------------------------------------

def slide_06_comms(prs):
    """Three-column section layout."""
    slide = blank_slide(prs)
    content_y = add_slide_header(slide,
                                 "Client Comms & QBR Readiness",
                                 slide_num=6)

    sections = [
        {
            "title": "Day-to-day principles",
            "items": [
                "— No surprises \u2014 flag risks early",
                "— Written confirmation of every verbal decision",
                "— Right channel: Slack (now) \u00b7 Email (record) \u00b7 Call (sensitive)",
            ],
        },
        {
            "title": "Launch support",
            "items": [
                "— \u221248h: pre-launch QA pack to Client Lead",
                "— Day-of: monitor first 2 hours live",
                "— +24h: short pacing & creative note to client",
            ],
        },
        {
            "title": "QBR \u2014 start 3 weeks out",
            "items": [
                "— Wk \u22123: data + narrative with Client Lead",
                "— Wk \u22122: draft + internal review",
                "— Wk \u22121: client pre-read + dry-run",
                "— Rule: zero numbers in the QBR the client hasn\u2019t already seen",
            ],
        },
    ]

    col_w  = Inches(3.8)
    col_xs = [Inches(0.5), Inches(4.7), Inches(8.9)]

    for col_x, section in zip(col_xs, sections):
        # Section heading
        add_textbox(slide, section["title"],
                    x=col_x, y=content_y,
                    w=col_w, h=Inches(0.38),
                    font_size=14, bold=True, color=BLACK)

        # Small accent rule under section heading
        add_accent_rule(slide,
                        x=col_x, y=content_y + Inches(0.42),
                        width=Inches(1.2))

        # Bullet items
        txBox = slide.shapes.add_textbox(
            col_x, content_y + Inches(0.6), col_w, Inches(4.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(section["items"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(8)
            run = p.add_run()
            run.text = item
            _fmt_run(run, FONT, 13, False, False, BLACK)


# ---------------------------------------------------------------------------
# Slide 7 — A Week in My Life
# ---------------------------------------------------------------------------

def slide_07_week(prs):
    """5-row Mon–Fri × AM/PM table + always-on footnote."""
    slide = blank_slide(prs)
    content_y = add_slide_header(slide,
                                 "A Week in My Life",
                                 slide_num=7,
                                 subtitle="on this account")

    # Table: header + 5 day rows
    table = slide.shapes.add_table(
        6, 3, Inches(0.5), content_y, Inches(12.3), Inches(4.8)
    ).table
    table.columns[0].width = Inches(1.7)
    table.columns[1].width = Inches(5.3)
    table.columns[2].width = Inches(5.3)

    for ci, hdr in enumerate(["Day", "AM", "PM"]):
        set_cell_text(table.cell(0, ci), hdr, font_size=13, bold=True, color=WHITE)

    week_data = [
        ("Mon", "Tracker review, stand-up",       "Weekly report build + QA"),
        ("Tue", "Send weekly report + actions",   "Bid review, search-term harvest"),
        ("Wed", "Client call w/ Client Lead",     "Builds + creative QA"),
        ("Thu", "Pacing & reallocations",         "Ad-hoc queries, QBR prep"),
        ("Fri", "Launches go live + monitoring",  "Friday wrap, plan next week"),
    ]
    for ri, (day, am, pm) in enumerate(week_data):
        set_cell_text(table.cell(ri + 1, 0), day, font_size=12, bold=True)
        set_cell_text(table.cell(ri + 1, 1), am,  font_size=12)
        set_cell_text(table.cell(ri + 1, 2), pm,  font_size=12)

    style_header_row(table)
    style_body_rows(table)

    # Always-on footnote in turquoise
    add_textbox(slide,
                "Always-on: issue triage  \u00b7  Slack monitoring  \u00b7  feed-health alerts",
                x=Inches(0.5), y=Inches(6.8),
                w=Inches(11), h=Inches(0.28),
                font_size=12, bold=True, color=TURQ)


# ---------------------------------------------------------------------------
# Slide 8 — Why I'm the Right CSA for Epsilon
# ---------------------------------------------------------------------------

def slide_08_why(prs):
    """6 mixed-run bullets + callout quote + Thank you."""
    slide = blank_slide(prs)
    content_y = add_slide_header(slide,
                                 "Why I\u2019m the Right CSA for Epsilon",
                                 slide_num=8)

    # Bullets: bold keyword + regular description
    bullets = [
        ("Organised & accountable",
         "single-source tracker, Friday wrap, no surprises"),
        ("Accurate",
         "6-point QA checklist, peer review on everything client-facing"),
        ("Clear communicator",
         "root cause + options + recommendation, every time"),
        ("Data-confident",
         "every chart has a \u201cso what?\u201d"),
        ("Collaborative",
         "make the Client Lead\u2019s job easier, the client\u2019s life simpler"),
        ("Excited about Epsilon",
         "CORE ID + closed-loop measurement + retailer\u2194brand seat"),
    ]

    txBox = slide.shapes.add_textbox(
        Inches(0.5), content_y, Inches(12.3), Inches(4.6)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (kw, desc) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        # Bold keyword
        r1 = p.add_run()
        r1.text = f"— {kw}"
        _fmt_run(r1, FONT, 14, True, False, BLACK)
        # Regular description
        r2 = p.add_run()
        r2.text = f"  \u2014  {desc}"
        _fmt_run(r2, FONT, 14, False, False, BLACK)

    # Closing callout quote
    add_callout(slide,
                "\u201cHappy to dig into any slide"
                " \u2014 particularly triage or optimisations.\u201d",
                x=Inches(0.5), y=Inches(6.1),
                w=Inches(10.5), h=Inches(0.55))

    # Thank you — right-aligned
    add_textbox(slide, "Thank you.",
                x=Inches(10.8), y=Inches(6.1),
                w=Inches(2.0), h=Inches(0.55),
                font_size=16, bold=True, color=BLACK,
                align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main():
    prs = new_prs()

    slide_01_title(prs)
    slide_02_operations(prs)
    slide_03_reporting(prs)
    slide_04_triage(prs)
    slide_05_optimisation(prs)
    slide_06_comms(prs)
    slide_07_week(prs)
    slide_08_why(prs)

    output = "Epsilon_CSA_Deck.pptx"
    prs.save(output)
    print(f"Saved: {output}")


if __name__ == "__main__":
    main()
